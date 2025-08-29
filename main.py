from __future__ import annotations
import io
import json
import os
import re
import sys
import time
from dataclasses import dataclass, field
from typing import Any, Dict, List, Optional

os.system("pip install streamlit==1.38.0")
os.system("pip install python-pptx==0.6.23")
os.system("pip install requests>=2.31.0")

import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.shapes import PP_PLACEHOLDER

import requests

APP_TITLE = "Your Text, Your Style – Auto-Generate a Presentation"
MAX_UPLOAD_MB = 20
MAX_TOKENS_FALLBACK = 8000  # naïve guardrail for giant inputs


# ------------------------------

def _write_repo_files_once():
    try:
        if not os.path.exists("LICENSE"):
            with open("LICENSE", "w", encoding="utf-8") as f:
                f.write(MIT_LICENSE_TEXT)
        if not os.path.exists("README.md"):
            with open("README.md", "w", encoding="utf-8") as f:
                f.write(README_MD)
        if not os.path.exists("requirements.txt"):
            with open("requirements.txt", "w", encoding="utf-8") as f:
                f.write(REQUIREMENTS_TXT)
    except Exception:
        pass

# ------------------------------
# LLM Provider Adapters
# ------------------------------
@dataclass
class LLMResponse:
    slides: List[Dict[str, Any]]


def call_openai(api_key: str, model: str, prompt: str, max_tokens: int = 2000) -> str:
    url = "https://api.openai.com/v1/chat/completions"
    headers = {"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"}
    body = {
        "model": model or "gpt-4o-mini",
        "messages": [
            {"role": "system", "content": "You are a helpful assistant that returns STRICT JSON only."},
            {"role": "user", "content": prompt},
        ],
        "response_format": {"type": "json_object"},
        "temperature": 0.2,
        "max_tokens": max_tokens,
    }
    r = requests.post(url, headers=headers, json=body, timeout=60)
    r.raise_for_status()
    return r.json()["choices"][0]["message"]["content"].strip()


def call_anthropic(api_key: str, model: str, prompt: str, max_tokens: int = 2000) -> str:
    url = "https://api.anthropic.com/v1/messages"
    headers = {"x-api-key": api_key, "anthropic-version": "2023-06-01", "content-type": "application/json"}
    body = {
        "model": model or "claude-3-5-sonnet-latest",
        "max_tokens": max_tokens,
        "temperature": 0.2,
        "messages": [{"role": "user", "content": prompt}],
        "system": "Return STRICT JSON only.",
    }
    r = requests.post(url, headers=headers, json=body, timeout=60)
    r.raise_for_status()
    parts = [c.get("text", "") for c in r.json().get("content", []) if c.get("type") == "text"]
    return "".join(parts).strip()


def call_gemini(api_key: str, model: str, prompt: str, max_tokens: int = 2048) -> str:
    model_name = model or "gemini-1.5-flash"
    url = f"https://generativelanguage.googleapis.com/v1beta/models/{model_name}:generateContent?key={api_key}"
    headers = {"content-type": "application/json"}
    body = {"contents": [{"parts": [{"text": prompt}]}],
            "generationConfig": {"temperature": 0.2, "maxOutputTokens": max_tokens}}
    r = requests.post(url, headers=headers, json=body, timeout=60)
    r.raise_for_status()
    data = r.json()
    try:
        return data["candidates"][0]["content"]["parts"][0]["text"].strip()
    except Exception as e:
        raise RuntimeError(f"Gemini response parse error: {e}\nRaw: {data}")


def outline_with_llm(provider: str, api_key: str, model: str, user_text: str, guidance: str, want_notes: bool) -> LLMResponse:
    schema_hint = {"slides": [{"title": "string", "bullets": ["string", "string"], **({"notes": "string"} if want_notes else {})}]}
    prompt = f"""
You are given INPUT_TEXT and optional GUIDANCE.
Return STRICT JSON matching this schema:\n{json.dumps(schema_hint)}

Rules:
- Choose a reasonable number of slides.
- Bullets concise (<=15 words).
- Headings → slide sections when useful.
- {"Include 'notes': expand beyond bullets with context/examples. Not duplication." if want_notes else "Do NOT include 'notes'."}

GUIDANCE: {guidance or "(none)"}
INPUT_TEXT:\n\n{user_text[:MAX_TOKENS_FALLBACK]}
"""
    try:
        if provider == "OpenAI":
            raw = call_openai(api_key, model, prompt)
        elif provider == "Anthropic":
            raw = call_anthropic(api_key, model, prompt)
        elif provider == "Gemini":
            raw = call_gemini(api_key, model, prompt)
        else:
            raise ValueError("Unknown provider")
        data = json.loads(raw)
    except Exception:
        return heuristic_outline(user_text, want_notes)

    slides = data.get("slides", [])
    norm_slides = []
    for s in slides:
        title = str(s.get("title", "Untitled")).strip()
        bullets = [str(b).strip() for b in s.get("bullets", []) if str(b).strip()]
        entry = {"title": title, "bullets": bullets}
        if want_notes and isinstance(s.get("notes"), str):
            entry["notes"] = s["notes"].strip()
        norm_slides.append(entry)
    return LLMResponse(slides=norm_slides)


# ------------------------------
# Non-LLM fallback slide mapping
# ------------------------------
MD_H1 = re.compile(r"^# +(.+)$", re.MULTILINE)
MD_H2 = re.compile(r"^## +(.+)$", re.MULTILINE)

def heuristic_outline(text: str, want_notes: bool) -> LLMResponse:
    text = text.strip()
    slides: List[Dict[str, Any]] = []
    headings = MD_H1.findall(text) or MD_H2.findall(text)
    if headings:
        sections = re.split(r"^##? +.+$", text, flags=re.MULTILINE)
        for i, section in enumerate(sections):
            if not section.strip():
                continue
            title = headings[i] if i < len(headings) else f"Section {i+1}"
            bullets = [line.strip("- *\t ") for line in section.splitlines() if line.strip().startswith(('-', '*'))]
            if not bullets:
                bullets = [s.strip() for s in re.split(r"(?<=[.!?])\s+", section) if len(s.strip()) > 0][:5]
            slides.append({"title": title.strip(), "bullets": bullets[:7], **({"notes": ""} if want_notes else {})})
    else:
        paras = [p.strip() for p in text.split('\n\n') if p.strip()]
        target = min(max(len(paras), 4), 10)
        chunks = max(1, len(paras) // target)
        for i in range(0, len(paras), chunks):
            blob = " ".join(paras[i:i+chunks])
            title = blob.split(". ")[0][:70] if blob else f"Slide {len(slides)+1}"
            bullets = [s.strip() for s in re.split(r"(?<=[.!?])\s+", blob) if s.strip()][:5]
            slides.append({"title": title or f"Slide {len(slides)+1}", "bullets": bullets[:7], **({"notes": ""} if want_notes else {})})
    if not slides:
        slides = [{"title": "Overview", "bullets": [text[:120]]}]
    return LLMResponse(slides=slides)

# ------------------------------
# Template analysis & PPT building
# ------------------------------
@dataclass
class TemplateAssets:
    picture_blobs: List[bytes] = field(default_factory=list)
    title_layout_idx: int = 0
    title_and_content_layout_idx: int = 0

def analyze_template(prs: Presentation) -> TemplateAssets:
    pics: List[bytes] = []
    for s in prs.slides:
        for shp in s.shapes:
            if shp.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    pics.append(shp.image.blob)
                except Exception:
                    pass
    title_idx, tac_idx = 0, 1 if len(prs.slide_layouts) > 1 else 0
    for i, layout in enumerate(prs.slide_layouts):
        names = ",".join([ph.name.lower() for ph in layout.placeholders])
        lname = (getattr(layout, 'name', '') or '').lower()
        if ("title" in names or "title" in lname) and ("content" not in names):
            title_idx = i
        if ("title" in names and ("content" in names or "body" in names)):
            tac_idx = i
    return TemplateAssets(picture_blobs=pics, title_layout_idx=title_idx, title_and_content_layout_idx=tac_idx)

def add_title_slide(prs: Presentation, assets: TemplateAssets, title: str, subtitle: str = ""):
    layout = prs.slide_layouts[assets.title_layout_idx]
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = title
        slide.shapes.title.text_frame.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.LEFT
    for ph in slide.placeholders:
        if ph.is_placeholder and 'subtitle' in (ph.name or '').lower():
            ph.text = subtitle
            break
    _maybe_add_reused_picture(slide, assets)

def _maybe_add_reused_picture(slide, assets: TemplateAssets):
    if not assets.picture_blobs:
        return
    blob = assets.picture_blobs[hash(slide) % len(assets.picture_blobs)]
    try:
        left = slide.part.slide_width - Inches(3.0)
        top = slide.part.slide_height - Inches(2.2)
        slide.shapes.add_picture(io.BytesIO(blob), left, top, width=Inches(2.8))
    except Exception:
        pass

def add_content_slide(prs: Presentation, assets: TemplateAssets, title: str, bullets: List[str], notes: Optional[str] = None):
    layout = prs.slide_layouts[assets.title_and_content_layout_idx]
    slide = prs.slides.add_slide(layout)
    if slide.shapes.title:
        slide.shapes.title.text = title

    body = None
    for ph in slide.placeholders:
        if ph.is_placeholder and ph.placeholder_format.type == PP_PLACEHOLDER.BODY:
            body = ph.text_frame
            break
    if body is None:
        left, top, width, height = Inches(1.0), Inches(1.8), Inches(8.0), Inches(4.5)
        body = slide.shapes.add_textbox(left, top, width, height).text_frame
    else:
        body.clear()

    for i, b in enumerate(bullets):
        p = body.add_paragraph() if i > 0 else body.paragraphs[0]
        p.text = b
        p.level = 0
        try:
            p.font.size = p.font.size or Pt(18)
        except Exception:
            pass

    if notes is not None:
        try:
            slide.notes_slide.notes_text_frame.text = notes
        except Exception:
            try:
                _ = slide.notes_slide
                slide.notes_slide.notes_text_frame.text = notes
            except Exception:
                pass

    _maybe_add_reused_picture(slide, assets)

# ------------------------------
# Streamlit UI
# ------------------------------

def _print_requirements_and_exit():
    print(REQUIREMENTS_TXT)
    sys.exit(0)

def main():
    if len(sys.argv) > 1 and sys.argv[1] == "--print-reqs":
        _print_requirements_and_exit()

    _write_repo_files_once()
    st.set_page_config(page_title=APP_TITLE, layout="wide")
    st.title(APP_TITLE)
    st.caption("Paste text → upload template → pick provider → download PPT. Keys never stored.")

    with st.sidebar:
        st.subheader("LLM Provider")
        provider = st.selectbox("Choose provider", ["None (Heuristic)", "OpenAI", "Anthropic", "Gemini"], index=0)
        model = st.text_input("Model (optional)")
        api_key = st.text_input("API Key", type="password") if provider != "None (Heuristic)" else ""
        want_notes = st.checkbox("Auto-generate speaker notes", value=True)
        preset = st.selectbox("Tone preset (optional)", ["", "Professional", "Investor pitch", "Visual-heavy", "Technical", "Educator"], index=0)

    col1, col2 = st.columns([2, 1])
    with col1:
        user_text = st.text_area("Paste your text or markdown", height=320, placeholder="# Title\n\nPaste lots of text…")
        guidance = st.text_input("One-line guidance (optional)", value=(preset or ""))
    with col2:
        uploaded = st.file_uploader("Upload PowerPoint template (.pptx/.potx)", type=["pptx", "potx"], accept_multiple_files=False)
        st.caption("We reuse its layouts, colors, and images.")
        gen_btn = st.button("Generate Presentation", type="primary", use_container_width=True)

    if gen_btn:
        if not user_text or not uploaded:
            st.error("Please provide text and a PowerPoint template.")
            return
        size_mb = (uploaded.size or 0) / (1024 * 1024)
        if size_mb > MAX_UPLOAD_MB:
            st.error(f"Template file too large ({size_mb:.1f} MB). Limit: {MAX_UPLOAD_MB} MB.")
            return

        try:
            prs = Presentation(uploaded)
            assets = analyze_template(prs)
            if provider == "None (Heuristic)":
                outline = heuristic_outline(user_text, want_notes)
            else:
                if not api_key:
                    st.error("Enter your API key or use heuristic mode.")
                    return
                outline = outline_with_llm(provider, api_key, model, user_text, guidance, want_notes)
        except Exception as e:
            st.error(f"Outline generation failed: {e}")
            return

        with st.expander("Preview Slides (Outline)"):
            for i, s in enumerate(outline.slides, 1):
                st.markdown(f"### Slide {i}: {s['title']}")
                st.write("• " + "\n• ".join(s['bullets']))
                if "notes" in s and s["notes"]:
                    st.info(f"Notes: {s['notes']}")

        try:
            template_bytes = uploaded.getvalue()
            prs_out = Presentation(io.BytesIO(template_bytes))
            if outline.slides:
                first = outline.slides[0]
                add_title_slide(prs_out, assets, first.get("title", ""), guidance or "")
                remaining = outline.slides[1:] if len(first.get("bullets", [])) <= 1 else outline.slides
            else:
                remaining = []
            for s in remaining:
                add_content_slide(prs_out, assets, s.get("title", "Untitled"), s.get("bullets", []), s.get("notes"))
            bio = io.BytesIO()
            prs_out.save(bio)
            bio.seek(0)
        except Exception as e:
            st.error(f"Failed to compose presentation: {e}")
            return

        st.success("Presentation ready!")
        st.download_button("Download .pptx", bio, "generated_presentation.pptx",
                           mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                           use_container_width=True)
        with st.expander("Show JSON outline (debug)"):
            st.json({"slides": outline.slides})

    st.markdown("---")
    st.markdown("**Privacy note:** API keys are never logged or stored.")

if __name__ == "__main__":
    main()
