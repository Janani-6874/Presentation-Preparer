from __future__ import annotations
import io
import json
import os
import re
import sys
import time
from dataclasses import dataclass, field
from typing import Any, Dict, List, Optional, Tuple


import streamlit as st
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.enum.shapes import MSO_SHAPE_TYPE


import requests


APP_TITLE = "Your Text, Your Style – Auto-Generate a Presentation"
MAX_UPLOAD_MB = 20
MAX_TOKENS_FALLBACK = 8000


REQUIREMENTS_TXT = """
streamlit==1.38.0
python-pptx==0.6.23
requests>=2.31.0
""".strip()




def _write_repo_files_once():
try:
if not os.path.exists("LICENSE"):
with open("LICENSE", "w", encoding="utf-8") as f:
f.write(MIT_LICENSE_TEXT)
if not os.path.exists("README.md"):
with open("README.md", "w", encoding="utf-8") as f:
f.write(README_MD)
if not os.path.exists("requirements.txt"):
with open("requirements.txt", "w", encoding="utf-8") as f:
f.write(REQUIREMENTS_TXT)
except Exception:
# Non-fatal; don't surface in UI to avoid noisy errors
pass


# Build outline
try:
if provider == "None (Heuristic)":
outline = heuristic_outline(user_text, want_notes)
else:
if not api_key:
st.error("Enter your API key for the selected provider, or choose the heuristic mode.")
return
outline = outline_with_llm(provider, api_key, model, user_text, guidance, want_notes)
except Exception as e:
st.error(f"Outline generation failed: {e}")
return


# LLM Provider Adapters
"temperature": 0.2,
"messages": [
{"role": "user", "content": prompt},
],
"system": "Return STRICT JSON only.",
}
r = requests.post(url, headers=headers, json=body, timeout=60)
r.raise_for_status()
data = r.json()
# Anthropic returns a list of content blocks; join text
parts = [c.get("text", "") for c in data.get("content", []) if c.get("type") == "text"]
return "".join(parts).strip()




def call_gemini(api_key: str, model: str, prompt: str, max_tokens: int = 2048) -> str:
# Google Generative Language API (Gemini 1.5) – JSON may vary per account/region
# Using the text endpoint with a JSON constraint in the prompt.
model_name = model or "gemini-1.5-flash"
url = f"https://generativelanguage.googleapis.com/v1beta/models/{model_name}:generateContent?key={api_key}"
headers = {"content-type": "application/json"}
body = {
"contents": [{"parts": [{"text": prompt}]}],
"generationConfig": {
"temperature": 0.2,
"maxOutputTokens": max_tokens,
},
}
r = requests.post(url, headers=headers, json=body, timeout=60)
r.raise_for_status()
data = r.json()
try:
return data["candidates"][0]["content"]["parts"][0]["text"].strip()
except Exception as e:
raise RuntimeError(f"Gemini response parse error: {e}\nRaw: {data}")




def outline_with_llm(provider: str, api_key: str, model: str, user_text: str, guidance: str, want_notes: bool) -> LLMResponse:
schema_hint = {
"slides": [
{
"title": "string",
"bullets": ["string", "string"],
**({"notes": "string"} if want_notes else {}),
}
]
}

prompt = f"""
if size_mb > MAX_UPLOAD_MB:
st.error(f"Template file too large ({size_mb:.1f} MB). Limit: {MAX_UPLOAD_MB} MB.")
return


with st.spinner("Analyzing template & mapping slides…"):
try:
prs = Presentation(uploaded)
except Exception as e:
st.error(f"Failed to read PowerPoint file: {e}")
return
assets = analyze_template(prs)


# Build outline
try:
if provider == "None (Heuristic)":
outline = heuristic_outline(user_text, want_notes)
else:
if not api_key:
st.error("Enter your API key for the selected provider, or choose the heuristic mode.")
return
outline = outline_with_llm(provider, api_key, model, user_text, guidance, want_notes)
except Exception as e:
st.error(f"Outline generation failed: {e}")
return


# Build deck
with st.spinner("Composing slides in your template style…"):
try:
# Start a fresh Presentation based on the same template to keep masters/theme
template_bytes = uploaded.getvalue()
prs_out = Presentation(io.BytesIO(template_bytes))


# Optionally add a title slide from first outline entry if it looks like a title
if outline.slides:
first = outline.slides[0]
add_title_slide(prs_out, assets, first.get("title", ""), guidance or "")
# Remove first from content if it’s only a title
remaining = outline.slides[1:] if (len(first.get("bullets", [])) <= 1) else outline.slides
else:
remaining = []


for s in remaining:
add_content_slide(prs_out, assets, s.get("title", "Untitled"), s.get("bullets", []), s.get("notes"))


bio = io.BytesIO()
prs_out.save(bio)
bio.seek(0)
except Exception as e:
st.error(f"Failed to compose presentation: {e}")
return


st.success("Presentation ready!")
st.download_button(
label="Download .pptx",
data=bio,
file_name="generated_presentation.pptx",
mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
use_container_width=True,
)


with st.expander("Show JSON outline (debug)"):
st.json({"slides": outline.slides})


st.markdown("---")
st.markdown(
"**Privacy note:** Your API key (if provided) is used only to call the selected provider from your browser session to this server. It is never logged or stored."
)




if __name__ == "__main__":
main()
